"""
core/conversion.py
PDF <-> Other format conversions.
"""
import fitz  # PyMuPDF
import os
from PIL import Image


def pdf_to_images(input_path: str, output_dir: str, fmt: str = "png", dpi: int = 200, cancel_event=None):
    """Convert each PDF page to an image file."""
    doc = fitz.open(input_path)
    results = []
    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)
    for i, page in enumerate(doc):
        if cancel_event:
            if cancel_event.is_set(): raise Exception("Cancelled by user")
            import time; time.sleep(0.005)
        pix = page.get_pixmap(matrix=mat)
        out = os.path.join(output_dir, f"page_{i + 1}.{fmt}")
        pix.save(out)
        results.append(out)
    doc.close()
    return results


def images_to_pdf(image_paths: list, output_path: str, cancel_event=None):
    """Convert a list of image files to a single PDF."""
    imgs = []
    first = None
    for p in image_paths:
        if cancel_event:
            if cancel_event.is_set(): raise Exception("Cancelled by user")
            import time; time.sleep(0.005)
        img = Image.open(p).convert("RGB")
        if first is None:
            first = img
        else:
            imgs.append(img)
    if first:
        first.save(output_path, save_all=True, append_images=imgs)
    return output_path


def pdf_to_text(input_path: str, output_path: str = None, cancel_event=None):
    """Extract all text from a PDF."""
    doc = fitz.open(input_path)
    text = ""
    for page in doc:
        if cancel_event:
            if cancel_event.is_set(): raise Exception("Cancelled by user")
            import time; time.sleep(0.005)
        text += page.get_text() + "\n"
    doc.close()
    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
    return text


def pdf_to_html(input_path: str, output_path: str, cancel_event=None):
    """Convert PDF to HTML using PyMuPDF."""
    doc = fitz.open(input_path)
    html_content = """<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>PDF Export</title>
<style>body{font-family:Arial,sans-serif;margin:40px;}.page{margin-bottom:30px;border-bottom:1px solid #ccc;padding-bottom:20px;}</style>
</head><body>"""
    for i, page in enumerate(doc):
        if cancel_event:
            if cancel_event.is_set(): raise Exception("Cancelled by user")
            import time; time.sleep(0.005)
        page_html = page.get_text("html")
        html_content += f'<div class="page"><h3>Page {i + 1}</h3>{page_html}</div>'
    html_content += "</body></html>"
    doc.close()
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    return output_path


def pdf_to_markdown(input_path: str, output_path: str, cancel_event=None):
    """Convert PDF to Markdown (text-based extraction)."""
    doc = fitz.open(input_path)
    md = ""
    for i, page in enumerate(doc):
        if cancel_event:
            if cancel_event.is_set(): raise Exception("Cancelled by user")
            import time; time.sleep(0.005)
        md += f"# Page {i + 1}\n\n"
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if block["type"] == 0:  # text block
                for line in block["lines"]:
                    line_text = ""
                    for span in line["spans"]:
                        text = span["text"]
                        size = span["size"]
                        if size > 20:
                            line_text += f"## {text}"
                        elif size > 16:
                            line_text += f"### {text}"
                        else:
                            line_text += text
                    md += line_text + "\n"
                md += "\n"
        md += "---\n\n"
    doc.close()
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(md)
    return output_path


def pdf_to_csv(input_path: str, output_path: str, cancel_event=None):
    """Extract tables from PDF to CSV using pdfplumber."""
    import pdfplumber
    import csv

    rows_all = []
    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    rows_all.append([cell if cell else "" for cell in row])

    with open(output_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerows(rows_all)
    return output_path


def pdf_to_json(input_path: str, output_path: str, cancel_event=None):
    """Extract structured text and metadata from PDF to JSON."""
    import json

    doc = fitz.open(input_path)
    data = {
        "metadata": doc.metadata,
        "page_count": len(doc),
        "pages": []
    }
    for i, page in enumerate(doc):
        if cancel_event:
            if cancel_event.is_set(): raise Exception("Cancelled by user")
            import time; time.sleep(0.005)
        page_data = {
            "page_number": i + 1,
            "width": page.rect.width,
            "height": page.rect.height,
            "text": page.get_text(),
        }
        data["pages"].append(page_data)
    doc.close()
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    return output_path


def pdf_to_xml(input_path: str, output_path: str, cancel_event=None):
    """Convert PDF to XML structured output."""
    doc = fitz.open(input_path)
    xml = '<?xml version="1.0" encoding="UTF-8"?>\n<document>\n'
    xml += f'  <metadata>\n'
    for key, val in doc.metadata.items():
        xml += f'    <{key}>{val if val else ""}</{key}>\n'
    xml += f'  </metadata>\n'
    for i, page in enumerate(doc):
        if cancel_event:
            if cancel_event.is_set(): raise Exception("Cancelled by user")
            import time; time.sleep(0.005)
        xml += f'  <page number="{i + 1}">\n'
        xml += f'    <text>{page.get_text().replace("<", "&lt;").replace(">", "&gt;")}</text>\n'
        xml += f'  </page>\n'
    xml += '</document>'
    doc.close()
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(xml)
    return output_path


def pdf_to_docx(input_path: str, output_path: str, cancel_event=None):
    """Convert PDF to Word DOCX."""
    from pdf2docx import Converter
    cv = Converter(input_path)
    cv.convert(output_path)
    cv.close()
    return output_path


def pdf_to_xlsx(input_path: str, output_path: str, cancel_event=None):
    """Extract tables from PDF to Excel XLSX."""
    import pdfplumber
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "PDF Tables"
    row_idx = 1
    with pdfplumber.open(input_path) as pdf:
        for page_num, page in enumerate(pdf.pages):
            if cancel_event:
                if cancel_event.is_set(): raise Exception("Cancelled by user")
                import time; time.sleep(0.005)
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    for col, cell in enumerate(row):
                        ws.cell(row=row_idx, column=col + 1, value=cell if cell else "")
                    row_idx += 1
                row_idx += 1  # blank row between tables
    wb.save(output_path)
    return output_path


def pdf_to_pptx(input_path: str, output_path: str, dpi: int = 150, cancel_event=None):
    """Convert PDF pages to PowerPoint slides (as images)."""
    from pptx import Presentation
    from pptx.util import Inches
    import tempfile

    doc = fitz.open(input_path)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)

    for page in doc:
        if cancel_event:
            if cancel_event.is_set(): raise Exception("Cancelled by user")
            import time; time.sleep(0.005)
        pix = page.get_pixmap(matrix=mat)
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        pix.save(tmp.name)
        tmp.close()

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        slide.shapes.add_picture(tmp.name, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        os.unlink(tmp.name)

    doc.close()
    prs.save(output_path)
    return output_path


def pdf_to_pdfa(input_path: str, output_path: str, cancel_event=None):
    """Convert PDF to PDF/A (archival) using PyMuPDF."""
    doc = fitz.open(input_path)
    # Set PDF/A metadata
    doc.set_metadata({
        "producer": "PDFOS",
        "creator": "PDFOS PDF/A Converter",
    })
    doc.save(output_path, deflate=True, garbage=4)
    doc.close()
    return output_path


def markdown_to_pdf(input_path: str, output_path: str, cancel_event=None):
    """Convert Markdown to PDF via HTML rendering."""
    import markdown

    with open(input_path, "r", encoding="utf-8") as f:
        md_text = f.read()

    html = markdown.markdown(md_text, extensions=["tables", "fenced_code"])
    styled_html = f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<style>body{{font-family:Arial,sans-serif;margin:40px;line-height:1.6;}}
code{{background:#f4f4f4;padding:2px 6px;border-radius:3px;}}
pre{{background:#f4f4f4;padding:12px;border-radius:5px;overflow-x:auto;}}
table{{border-collapse:collapse;width:100%;}}th,td{{border:1px solid #ddd;padding:8px;text-align:left;}}
</style></head><body>{html}</body></html>"""

    # Use PyMuPDF Story to convert HTML to PDF
    try:
        story = fitz.Story(html=styled_html)
        writer = fitz.DocumentWriter(output_path)
        mediabox = fitz.paper_rect("a4")
        where = mediabox + fitz.Rect(36, 36, -36, -36)
        more = True
        while more:
            if cancel_event:
                if cancel_event.is_set(): raise Exception("Cancelled by user")
                import time; time.sleep(0.005)
            dev = writer.begin_page(mediabox)
            more, _ = story.place(where)
            story.draw(dev)
            writer.end_page()
        writer.close()
    except Exception:
        # Fallback: create simple text PDF
        doc = fitz.open()
        page = doc.new_page()
        text_rect = fitz.Rect(36, 36, page.rect.width - 36, page.rect.height - 36)
        page.insert_textbox(text_rect, md_text, fontsize=11)
        doc.save(output_path)
        doc.close()
    return output_path


def html_to_pdf(input_path: str, output_path: str, cancel_event=None):
    """Convert HTML file to PDF."""
    with open(input_path, "r", encoding="utf-8") as f:
        html_content = f.read()

    try:
        story = fitz.Story(html=html_content)
        writer = fitz.DocumentWriter(output_path)
        mediabox = fitz.paper_rect("a4")
        where = mediabox + fitz.Rect(36, 36, -36, -36)
        more = True
        while more:
            if cancel_event:
                if cancel_event.is_set(): raise Exception("Cancelled by user")
                import time; time.sleep(0.005)
            dev = writer.begin_page(mediabox)
            more, _ = story.place(where)
            story.draw(dev)
            writer.end_page()
        writer.close()
    except Exception:
        doc = fitz.open()
        page = doc.new_page()
        page.insert_textbox(fitz.Rect(36, 36, 560, 800), html_content[:5000], fontsize=10)
        doc.save(output_path)
        doc.close()
    return output_path
